import os
import io
import csv
import json
import xml.etree.ElementTree as ET
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
import openpyxl
import pandas as pd
import shutil
import chardet


def parse_pptx(file_path):
    presentation = Presentation(file_path)
    slides = []
    for i, slide in enumerate(presentation.slides, 1):
        slide_content = []
        for shape in slide.shapes:
            # If the slide has a text attribute, it will be parsed
            if hasattr(shape, 'text'):
                slide_content.append(shape.text)
        slides.append({f"Slide {i}": slide_content})
    return {"slides": slides}


def parse_doc(file_path):
    document = Document(file_path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs]
    return {"paragraphs": paragraphs}


def parse_pdf(file_path):
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages, 1):
        pages.append({f"Page {i}": page.extract_text()})
    return {"pages": pages}


def parse_xml(file_path):
    # This function converts an XML file to a dictionary
    tree = ET.parse(file_path)
    root = tree.getroot()

    def element_to_dict(element):
        result = {}
        for child in element:
            child_data = element_to_dict(child)
            if child.tag in result:
                if type(result[child.tag]) is list:
                    result[child.tag].append(child_data)
                else:
                    result[child.tag] = [result[child.tag], child_data]
            else:
                result[child.tag] = child_data
        if element.text and element.text.strip():
            if len(result) > 0:
                result['#text'] = element.text.strip()
            else:
                result = element.text.strip()
        return result

    return element_to_dict(root)


def parse_excel(file_path):
    workbook = openpyxl.load_workbook(file_path)
    sheets = {}
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        data = []
        # Current setup only parses the content of the cell and not:
        # - Coordinate: The cell's address (e.g., 'A1').
        # - Row: The row number of the cell.
        # - Column: The column number of the cell.
        # - Data Type: The type of data stored in the cell (e.g., string, number, formula).
        # - Style: The style applied to the cell (e.g., font, fill, border).
        # - Formula: If the cell contains a formula, the formula itself.
        # - Hyperlink: If the cell contains a hyperlink, the hyperlink itself.
        # - Comment: If the cell contains a comment, the comment itself.
        # - Number Format: The number format applied to the cell (e.g., currency, date).
        # - Protection: Information about cell protection (e.g., locked, hidden).
        # If this kind of data is required turn values_only to False
        for row in sheet.iter_rows(values_only=True):
            data.append(list(row))
        sheets[sheet_name] = data
    return {"sheets": sheets}


def parse_csv(file_path):
    # Current parser works on the assumption that we re dealing with polish signs thus utf-8
    with open(file_path, 'r', encoding='utf-8') as file:
        reader = csv.reader(file)
        data = list(reader)
    return {"data": data}


def detect_encoding(file_path):
    with open(file_path, 'rb') as file:
        raw = file.read()
    return chardet.detect(raw)['encoding']


def parse_txt(file_path):
    encoding = detect_encoding(file_path)
    with open(file_path, 'r', encoding=encoding) as file:
        content = file.read()
    return {"content": content}


def parse_rtf(file_path):
    # This is a simple approach and may not handle all RTF formatting
    encoding = detect_encoding(file_path)
    with open(file_path, 'r', encoding=encoding, errors='ignore') as file:
        content = file.read()
    return {"content": content}


def parse_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    # Types of supported files in our parser
    parsers = {
        '.pptx': parse_pptx,
        '.doc': parse_doc,
        '.docx': parse_doc,
        '.pdf': parse_pdf,
        '.xml': parse_xml,
        '.xlsx': parse_excel,
        '.xls': parse_excel,
        '.csv': parse_csv,
        '.txt': parse_txt,
        '.rtf': parse_rtf
    }

    parser = parsers.get(file_extension)
    if parser:
        return parser(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")


# Version without creating separate files in output folder
# def process_folder(input_folder, output_folder):
#     for root, dirs, files in os.walk(input_folder):
#         for file in files:
#             input_file_path = os.path.join(root, file)
#             relative_path = os.path.relpath(input_file_path, input_folder)
#             output_file_path = os.path.join(output_folder, relative_path)
#             output_file_path = os.path.splitext(output_file_path)[0] + '.json'
#
#             os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
#
#             try:
#                 parsed_content = parse_file(input_file_path)
#                 with open(output_file_path, 'w', encoding='utf-8') as f:
#                     json.dump(parsed_content, f, ensure_ascii=False, indent=2)
#                 print(f"Processed: {input_file_path} -> {output_file_path}")
#             except Exception as e:
#                 print(f"Error processing {input_file_path}: {str(e)}")


def process_folder(input_folder, output_folder):
    for root, dirs, files in os.walk(input_folder):
        for file in files:
            input_file_path = os.path.join(root, file)
            relative_path = os.path.relpath(input_file_path, input_folder)
            file_extension = os.path.splitext(file)[1].lower()
            extension_folder = os.path.join(output_folder, file_extension.lstrip('.'))
            output_file_path = os.path.join(extension_folder, relative_path)
            output_file_path = os.path.splitext(output_file_path)[0] + '.json'

            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)

            try:
                parsed_content = parse_file(input_file_path)
                with open(output_file_path, 'w', encoding='utf-8') as f:
                    json.dump(parsed_content, f, ensure_ascii=False, indent=2)
                print(f"Processed: {input_file_path} -> {output_file_path}")
            except Exception as e:
                print(f"Error processing {input_file_path}: {str(e)}")


def main():
    # Input folder should contain files to be parsed
    base_folder = r"C:\Users\szyme\to_be_parsed"
    input_folder = os.path.join(base_folder, "input")
    output_folder = os.path.join(base_folder, "output")

    if not os.path.exists(input_folder):
        os.makedirs(input_folder)
        print(f"Created input folder: {input_folder}")

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
        print(f"Created output folder: {output_folder}")

    process_folder(input_folder, output_folder)


if __name__ == "__main__":
    main()
