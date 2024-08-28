import os
import io
import csv
import json
import xml.etree.ElementTree as ET
from pptx import Presentation
from docx import Document
from PyPDF2 import PdfReader
import openpyxl
import chardet
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import subprocess

# Ensure you have Tesseract installed and set the path
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Ustaw ścieżkę do Popplera względem katalogu projektu
# POPPLER_PATH = os.path.join(os.path.dirname(__file__), 'poppler', 'Library', 'bin')
PDFTOCAIRO_PATH = r'C:\Users\szyme\PycharmProjects\Azure-db-test2\poppler\Library\bin'
print(f"Ścieżka do Popplera: {PDFTOCAIRO_PATH}")
os.environ["PATH"] += os.pathsep + PDFTOCAIRO_PATH


# def check_poppler_installation():
#     pdftocairo_path = PDFTOCAIRO_PATH
#     if os.path.exists(pdftocairo_path):
#         # print(f"Znaleziono pdftocairo.exe: {pdftocairo_path}")
#         return pdftocairo_path
#     else:
#         print("Nie znaleziono pdftocairo.exe")
#         return None

def check_poppler_installation():
    pdftocairo_path = PDFTOCAIRO_PATH
    if os.path.exists(os.path.join(pdftocairo_path, 'pdftocairo.exe')) and os.path.exists(
            os.path.join(pdftocairo_path, 'pdfinfo.exe')):
        return pdftocairo_path
    else:
        print("Poppler utilities (pdftocairo and pdfinfo) not found in the specified path.")
        return None


def parse_pptx(file_path):
    try:
        presentation = Presentation(file_path)
        slides = []
        for i, slide in enumerate(presentation.slides, 1):
            slide_content = []
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text'):
                        slide_content.append(shape.text)
                    elif shape.shape_type == 13:  # 13 to wartość dla obrazów
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            pil_image = Image.open(io.BytesIO(image_bytes))
                            ocr_text = pytesseract.image_to_string(pil_image)
                            slide_content.append(f"OCR text from image: {ocr_text}")
                        except Exception as img_error:
                            print(f"Nie można przetworzyć obrazu na slajdzie {i}: {str(img_error)}")
                            slide_content.append("[Nieobsługiwany format obrazu]")
                except Exception as e:
                    print(f"Błąd podczas przetwarzania elementu na slajdzie {i}: {str(e)}")
            slides.append({f"Slajd {i}": slide_content})
        return {"slajdy": slides}
    except Exception as e:
        print(f"Błąd podczas przetwarzania pliku PPTX {file_path}: {str(e)}")
        return {"error": str(e)}


def parse_doc(file_path):
    document = Document(file_path)
    paragraphs = [paragraph.text for paragraph in document.paragraphs]
    return {"paragraphs": paragraphs}


# def parse_pdf(file_path):
#     pdftocairo_path = check_poppler_installation()
#     if not pdftocairo_path:
#         return {"error": "Nie znaleziono wymaganego narzędzia pdftocairo.exe"}
#
#     reader = PdfReader(file_path)
#     pages = []
#     for i, page in enumerate(reader.pages, 1):
#         text = page.extract_text()
#         if not text.strip():  # Jeśli nie wyodrębniono tekstu, próbuj OCR
#             try:
#                 print(f"Próba konwersji strony {i} pliku PDF na obraz...")
#                 images = convert_from_path(file_path, first_page=i, last_page=i, poppler_path=PDFTOCAIRO_PATH)
#                 print(f"Konwersja udana. Liczba obrazów: {len(images)}")
#                 for image in images:
#                     text += pytesseract.image_to_string(image)
#             except Exception as e:
#                 print(f"Błąd podczas konwersji strony {i} pliku PDF na obraz: {str(e)}")
#                 print(f"Typ wyjątku: {type(e).__name__}")
#                 text += f"[Błąd OCR na stronie {i}]"
#         pages.append({f"Strona {i}": text})
#     return {"strony": pages}

def parse_pdf(file_path):
    poppler_path = check_poppler_installation()
    if not poppler_path:
        return {"error": "Poppler utilities not found."}

    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text()
        if not text or not text.strip():  # If no text was extracted, try using OCR
            try:
                print(f"Attempting to convert page {i} of the PDF to an image for OCR...")
                images = convert_from_path(file_path, first_page=i, last_page=i, poppler_path=poppler_path)
                print(f"Page {i} converted successfully to image(s).")
                for image in images:
                    ocr_text = pytesseract.image_to_string(image)
                    text += ocr_text
            except Exception as e:
                print(f"Error converting page {i} of PDF to image: {str(e)}")
                text += f"[OCR error on page {i}]"
        else:
            print(f"Text extracted from page {i} successfully.")
        pages.append({f"Page {i}": text})

    return {"pages": pages}


def parse_xml(file_path):
    # This function converts an XML file to a dictionary
    tree = ET.parse(file_path)
    root = tree.getroot()

    def element_to_dict(element):
        result = {}
        for child in element:
            child_data = element_to_dict(child)
            if child.tag in result:
                if type(result[child.tag]) is list:
                    result[child.tag].append(child_data)
                else:
                    result[child.tag] = [result[child.tag], child_data]
            else:
                result[child.tag] = child_data
        if element.text and element.text.strip():
            if len(result) > 0:
                result['#text'] = element.text.strip()
            else:
                result = element.text.strip()
        return result

    return element_to_dict(root)


def parse_excel(file_path):
    workbook = openpyxl.load_workbook(file_path)
    sheets = {}
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        data = []
        # Current setup only parses the content of the cell and not:
        # - Coordinate: The cell's address (e.g., 'A1').
        # - Row: The row number of the cell.
        # - Column: The column number of the cell.
        # - Data Type: The type of data stored in the cell (e.g., string, number, formula).
        # - Style: The style applied to the cell (e.g., font, fill, border).
        # - Formula: If the cell contains a formula, the formula itself.
        # - Hyperlink: If the cell contains a hyperlink, the hyperlink itself.
        # - Comment: If the cell contains a comment, the comment itself.
        # - Number Format: The number format applied to the cell (e.g., currency, date).
        # - Protection: Information about cell protection (e.g., locked, hidden).
        # If this kind of data is required turn values_only to False
        for row in sheet.iter_rows(values_only=True):
            data.append(list(row))
        sheets[sheet_name] = data
    return {"sheets": sheets}


def parse_csv(file_path):
    # Current parser works on the assumption that we re dealing with polish signs thus utf-8
    with open(file_path, 'r', encoding='utf-8') as file:
        reader = csv.reader(file)
        data = list(reader)
    return {"data": data}


def detect_encoding(file_path):
    with open(file_path, 'rb') as file:
        raw = file.read()
    return chardet.detect(raw)['encoding']


def parse_txt(file_path):
    encoding = detect_encoding(file_path)
    with open(file_path, 'r', encoding=encoding) as file:
        content = file.read()
    return {"content": content}


def parse_rtf(file_path):
    # This is a simple approach and may not handle all RTF formatting
    encoding = detect_encoding(file_path)
    with open(file_path, 'r', encoding=encoding, errors='ignore') as file:
        content = file.read()
    return {"content": content}


def parse_file(file_path):
    _, file_extension = os.path.splitext(file_path)
    file_extension = file_extension.lower()

    # Types of supported files in our parser
    parsers = {
        '.pptx': parse_pptx,
        '.doc': parse_doc,
        '.docx': parse_doc,
        '.pdf': parse_pdf,
        '.xml': parse_xml,
        '.xlsx': parse_excel,
        '.xls': parse_excel,
        '.csv': parse_csv,
        '.txt': parse_txt,
        '.rtf': parse_rtf
    }

    parser = parsers.get(file_extension)
    if parser:
        try:
            return parser(file_path)
        except Exception as e:
            print(f"Błąd podczas przetwarzania pliku {file_path}: {str(e)}")
            return {"error": str(e)}
    else:
        raise ValueError(f"Nieobsługiwany format pliku: {file_extension}")


# Version without creating separate files in output folder
# def process_folder(input_folder, output_folder):
#     for root, dirs, files in os.walk(input_folder):
#         for file in files:
#             input_file_path = os.path.join(root, file)
#             relative_path = os.path.relpath(input_file_path, input_folder)
#             output_file_path = os.path.join(output_folder, relative_path)
#             output_file_path = os.path.splitext(output_file_path)[0] + '.json'
#
#             os.makedirs(os.path.dirname(output_file_path), exist_ok=True)
#
#             try:
#                 parsed_content = parse_file(input_file_path)
#                 with open(output_file_path, 'w', encoding='utf-8') as f:
#                     json.dump(parsed_content, f, ensure_ascii=False, indent=2)
#                 print(f"Processed: {input_file_path} -> {output_file_path}")
#             except Exception as e:
#                 print(f"Error processing {input_file_path}: {str(e)}")


def process_folder(input_folder, output_folder):
    for root, dirs, files in os.walk(input_folder):
        for file in files:
            input_file_path = os.path.join(root, file)
            relative_path = os.path.relpath(input_file_path, input_folder)
            file_extension = os.path.splitext(file)[1].lower()
            extension_folder = os.path.join(output_folder, file_extension.lstrip('.'))
            output_file_path = os.path.join(extension_folder, relative_path)
            output_file_path = os.path.splitext(output_file_path)[0] + '.json'

            os.makedirs(os.path.dirname(output_file_path), exist_ok=True)

            try:
                parsed_content = parse_file(input_file_path)
                with open(output_file_path, 'w', encoding='utf-8') as f:
                    json.dump(parsed_content, f, ensure_ascii=False, indent=2)
                print(f"Processed: {input_file_path} -> {output_file_path}")
            except Exception as e:
                print(f"Error processing {input_file_path}: {str(e)}")


def main():
    # Input folder should contain files to be parsed
    base_folder = r"C:\Users\szyme\to_be_parsed"
    input_folder = os.path.join(base_folder, "input")
    output_folder = os.path.join(base_folder, "output")

    if not os.path.exists(input_folder):
        os.makedirs(input_folder)
        print(f"Created input folder: {input_folder}")

    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
        print(f"Created output folder: {output_folder}")

    process_folder(input_folder, output_folder)


if __name__ == "__main__":
    main()
